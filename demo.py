"""
提取pdf文字
"""
import json
import os

from pptx import Presentation
import PyPDF2


def extract_text_from_pdf(file_path):
    text = ''
    with open(file_path, 'rb') as file:
        pdf_reader = PyPDF2.PdfFileReader(file)
        num_pages = pdf_reader.getNumPages()

        for page_number in range(num_pages):
            page = pdf_reader.getPage(page_number)
            text += page.extractText()

    return text


def extract_text_from_ppt(pptx_file):
    presentation = Presentation(pptx_file)
    text_runs = []

    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_runs.append(shape.text)
            elif shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text_runs.append(run.text)

    return '\n'.join(text_runs)


# 使用方法

if __name__ == '__main__':
    with open('./config/config.json', 'r') as f:
        config = json.load(f)

    filepath = config['filepath']
    target = config['target']

    for filename in os.listdir(filepath):
        path = os.path.join(filepath, filename)
        abspath = os.path.abspath(path)
        name = filename.split('.')[0]
        if not os.path.exists(abspath):
            print('路径不存在')
        else:
            text = extract_text_from_pdf(abspath)
            bs = bytes(text, encoding='utf-8')
            target_name = target + name + '.txt'
            with open(target_name, 'wb') as f:
                f.write(bs)
                f.close()
